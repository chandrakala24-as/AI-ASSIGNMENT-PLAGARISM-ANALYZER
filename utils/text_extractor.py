"""
text_extractor.py
=================
Lightweight, Render-free-tier-compatible document text extraction.

Priority chain for every file type:
  PDF   → pypdf native text  →  (if <50 chars) page-image render + cloud OCR
  DOCX  → python-docx        →  xml fallback
  PPTX  → python-pptx        →  xml fallback
  Image → OCR.space cloud API (zero RAM, no torch/easyocr required)

OCR.space (https://ocr.space) offers a free API key ('helloworld') that
accepts base64-encoded images and returns parsed text. It runs in the cloud
so Render's 512 MB RAM limit is never a concern.
"""

import os
import io
import re
import gc
import base64
import json
import zipfile
import xml.etree.ElementTree as ET
import urllib.request
import urllib.parse
import urllib.error
from typing import Optional

import fitz  # PyMuPDF
from PIL import Image

# ─── optional library imports ────────────────────────────────────────────────

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import easyocr
    import numpy as np
    # Disable local EasyOCR on Render free tier or if requested, to avoid RAM OOM crashes.
    DISABLE_LOCAL_OCR = os.getenv("DISABLE_LOCAL_OCR", "false").lower() == "true" or os.getenv("RENDER") == "true"
    HAS_EASYOCR = not DISABLE_LOCAL_OCR
except ImportError:
    HAS_EASYOCR = False

# ─── OCR.space cloud OCR (primary OCR backend – zero local RAM needed) ────────

_OCR_SPACE_KEY = os.getenv("OCR_SPACE_API_KEY", "helloworld")  # free public key


def _ocr_space_image(pil_img: Image.Image, language: str = "eng") -> str:
    """
    Send a PIL image to OCR.space free API and return recognised text.
    Works on any Render/Vercel deployment — no GPU, no torch, <1 MB payload.
    """
    try:
        # Downscale large images to stay well under the 1 MB API limit
        img_copy = pil_img.copy().convert("RGB")
        img_copy.thumbnail((1400, 1400), Image.LANCZOS)

        buf = io.BytesIO()
        img_copy.save(buf, format="JPEG", quality=85, optimize=True)
        encoded = base64.b64encode(buf.getvalue()).decode("utf-8")

        payload = urllib.parse.urlencode({
            "apikey":      _OCR_SPACE_KEY,
            "base64image": f"data:image/jpeg;base64,{encoded}",
            "language":    language,
            "isTable":     "false",
            "scale":       "true",
            "OCREngine":   "2",          # Engine 2 is better for mixed / handwritten text
        }).encode("utf-8")

        req = urllib.request.Request(
            "https://api.ocr.space/parse/image",
            data=payload,
            headers={"User-Agent": "PlagCheckAI/1.0"},
        )

        with urllib.request.urlopen(req, timeout=20) as resp:
            result = json.loads(resp.read().decode("utf-8"))

        if result.get("IsErroredOnProcessing"):
            print(f"[OCR.space] API error: {result.get('ErrorMessage')}")
            return ""

        parts = result.get("ParsedResults") or []
        text = " ".join(p.get("ParsedText", "") for p in parts).strip()
        if text:
            print(f"[OCR.space] Extracted {len(text.split())} words via cloud OCR.")
        return text

    except urllib.error.URLError as net_err:
        print(f"[OCR.space] Network error (offline?): {net_err}")
        return ""
    except Exception as exc:
        print(f"[OCR.space] Unexpected error: {exc}")
        return ""


_EASYOCR_READER = None

def _ocr_image(pil_img: Image.Image, language: str = "eng") -> str:
    """
    Perform OCR on a PIL Image.
    First tries local EasyOCR if available and not explicitly disabled.
    Otherwise, falls back to OCR.space cloud API.
    """
    global _EASYOCR_READER
    if HAS_EASYOCR:
        try:
            if _EASYOCR_READER is None:
                print("[OCR] Initialising local EasyOCR reader (this may take a few seconds on first run)...")
                _EASYOCR_READER = easyocr.Reader(['en'], gpu=False)
            
            # Convert PIL Image to RGB NumPy array as expected by EasyOCR
            img_np = np.array(pil_img.convert("RGB"))
            result = _EASYOCR_READER.readtext(img_np, detail=0)
            text = " ".join(result).strip()
            if text:
                print(f"[OCR] Extracted {len(text.split())} words via local EasyOCR.")
                return text
        except Exception as exc:
            print(f"[OCR] Local EasyOCR failed: {exc}. Falling back to cloud OCR.")

    # Fallback to OCR.space cloud API
    return _ocr_space_image(pil_img, language)


# ─── PDF rendering for scanned pages (uses pypdf + Pillow, no extra deps) ────

def _render_pdf_page_to_pil(page) -> Optional[Image.Image]:
    """
    Render a PyMuPDF page to a PIL Image for OCR fallback.
    """
    try:
        pix = page.get_pixmap(dpi=150)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        return img
    except Exception as exc:
        print(f"[PDF] Could not render page image: {exc}")
    return None


# ─── PDF extraction ──────────────────────────────────────────────────────────

def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from a PDF.
    1. Native text layer (PyMuPDF) — fastest, no deps.
    2. If result is too short (<50 chars per page avg), assume scanned PDF:
       render pages to images, send each to OCR.space.
    """
    all_text: list[str] = []

    try:
        doc = fitz.open(file_path)
        for page in doc:
            page_text = page.get_text() or ""
            all_text.append(page_text)
    except Exception as exc:
        print(f"[PDF] PyMuPDF failed: {exc}")
        return ""

    combined = "\n".join(all_text).strip()

    # If native extraction produced meaningful text — done.
    if len(combined) >= 50:
        print(f"[PDF] Native extraction: {len(combined.split())} words.")
        return combined

    # ---------- Scanned / image-based PDF → cloud OCR ----------
    print(f"[PDF] Native text too short ({len(combined)} chars). Attempting OCR on page images…")
    ocr_parts: list[str] = []

    try:
        for page_num, page in enumerate(doc, 1):
            if page_num > 10:  # Limit OCR to first 10 pages to avoid timeouts on Render
                print(f"[PDF] Reached 10 page limit for OCR, stopping to prevent timeouts.")
                break
                
            pil_img = _render_pdf_page_to_pil(page)
            if pil_img is None:
                continue

            text = _ocr_image(pil_img)
            if text:
                ocr_parts.append(text)
            else:
                print(f"[PDF] Page {page_num}: OCR returned empty.")

            del pil_img
            gc.collect()

    except Exception as exc:
        print(f"[PDF] OCR pass failed: {exc}")

    if ocr_parts:
        final = "\n".join(ocr_parts).strip()
        print(f"[PDF] OCR total: {len(final.split())} words from {len(ocr_parts)} pages.")
        return final

    # Return whatever native text we have (even if short)
    return combined


# ─── DOCX extraction ─────────────────────────────────────────────────────────

def _extract_docx_xml_fallback(file_path: str) -> str:
    """Zero-dependency DOCX extraction via zip+XML."""
    try:
        with zipfile.ZipFile(file_path) as zf:
            xml_bytes = zf.read("word/document.xml")
        root = ET.fromstring(xml_bytes)
        NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        paragraphs = []
        for para in root.iter(f"{{{NS}}}p"):
            runs = "".join(
                t.text for t in para.iter(f"{{{NS}}}t") if t.text
            )
            if runs:
                paragraphs.append(runs)
        return "\n".join(paragraphs)
    except Exception as exc:
        print(f"[DOCX] XML fallback failed: {exc}")
        return ""


def extract_text_from_docx(file_path: str) -> str:
    if HAS_DOCX:
        try:
            doc = docx.Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs)
            if text.strip():
                return text
        except Exception as exc:
            print(f"[DOCX] python-docx failed: {exc}")
    return _extract_docx_xml_fallback(file_path)


# ─── PPTX extraction ─────────────────────────────────────────────────────────

def _extract_pptx_xml_fallback(file_path: str) -> str:
    try:
        parts: list[str] = []
        with zipfile.ZipFile(file_path) as zf:
            slides = sorted(
                n for n in zf.namelist()
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            )
            for slide_file in slides:
                root = ET.fromstring(zf.read(slide_file))
                NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
                slide_text = " ".join(
                    t.text.strip() for t in root.iter(f"{{{NS}}}t") if t.text and t.text.strip()
                )
                if slide_text:
                    parts.append(slide_text)
        return "\n".join(parts)
    except Exception as exc:
        print(f"[PPTX] XML fallback failed: {exc}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    if HAS_PPTX:
        try:
            prs = Presentation(file_path)
            parts: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if slide_texts:
                    parts.append(f"[Slide {i}] " + " ".join(slide_texts))
            if parts:
                return "\n".join(parts)
        except Exception as exc:
            print(f"[PPTX] python-pptx failed: {exc}")
    return _extract_pptx_xml_fallback(file_path)


# ─── Image extraction ─────────────────────────────────────────────────────────

def extract_text_from_image(file_path: str) -> str:
    """Send image directly to OCR.space cloud API."""
    try:
        pil_img = Image.open(file_path)
        text = _ocr_image(pil_img)
        return text
    except Exception as exc:
        print(f"[IMG] Failed to open image for OCR: {exc}")
        return ""


# ─── Public entry point ───────────────────────────────────────────────────────

def extract_text(file_path: str) -> str:
    """
    Dispatch to the correct extractor based on file extension.
    Returns an empty string on failure (never raises).
    """
    if not os.path.exists(file_path):
        print(f"[Extractor] File not found: {file_path}")
        return ""

    ext = os.path.splitext(file_path)[1].lower()
    print(f"[Extractor] Processing {os.path.basename(file_path)} ({ext})")

    try:
        if ext == ".pdf":
            return extract_text_from_pdf(file_path)
        elif ext == ".docx":
            return extract_text_from_docx(file_path)
        elif ext in (".pptx", ".ppt"):
            return extract_text_from_pptx(file_path)
        elif ext in (".png", ".jpg", ".jpeg"):
            return extract_text_from_image(file_path)
        else:
            # Generic text fallback
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception as exc:
        print(f"[Extractor] Unhandled error for {file_path}: {exc}")
        return ""
